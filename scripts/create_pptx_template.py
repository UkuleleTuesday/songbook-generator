#!/usr/bin/env python3
"""
Creates generator/templates/song_slide_template.pptx from scratch.

Run once with: uv run scripts/create_pptx_template.py
"""
import io
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu

# Exact EMU values measured from Love_Me_Do.pptx
SLIDE_W = 13_716_000
SLIDE_H = 10_287_000

GROUP_XML = """\
<p:grpSp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
         xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:nvGrpSpPr>
    <p:cNvPr name="Group 2" id="2"/>
    <p:cNvGrpSpPr/>
    <p:nvPr/>
  </p:nvGrpSpPr>
  <p:grpSpPr>
    <a:xfrm rot="0">
      <a:off x="1714500" y="0"/>
      <a:ext cx="10287000" cy="10287000"/>
      <a:chOff x="0" y="0"/>
      <a:chExt cx="812800" cy="812800"/>
    </a:xfrm>
  </p:grpSpPr>
  <p:sp>
    <p:nvSpPr>
      <p:cNvPr name="Freeform 3" id="3"/>
      <p:cNvSpPr/>
      <p:nvPr/>
    </p:nvSpPr>
    <p:spPr>
      <a:xfrm flipH="false" flipV="false" rot="0">
        <a:off x="0" y="0"/>
        <a:ext cx="812800" cy="812800"/>
      </a:xfrm>
      <a:custGeom>
        <a:avLst/>
        <a:gdLst/>
        <a:ahLst/>
        <a:cxnLst/>
        <a:rect r="r" b="b" t="t" l="l"/>
        <a:pathLst>
          <a:path h="812800" w="812800">
            <a:moveTo><a:pt x="406400" y="0"/></a:moveTo>
            <a:cubicBezTo>
              <a:pt x="181951" y="0"/>
              <a:pt x="0" y="181951"/>
              <a:pt x="0" y="406400"/>
            </a:cubicBezTo>
            <a:cubicBezTo>
              <a:pt x="0" y="630849"/>
              <a:pt x="181951" y="812800"/>
              <a:pt x="406400" y="812800"/>
            </a:cubicBezTo>
            <a:cubicBezTo>
              <a:pt x="630849" y="812800"/>
              <a:pt x="812800" y="630849"/>
              <a:pt x="812800" y="406400"/>
            </a:cubicBezTo>
            <a:cubicBezTo>
              <a:pt x="812800" y="181951"/>
              <a:pt x="630849" y="0"/>
              <a:pt x="406400" y="0"/>
            </a:cubicBezTo>
            <a:close/>
          </a:path>
        </a:pathLst>
      </a:custGeom>
      <a:solidFill>
        <a:srgbClr val="000000"/>
      </a:solidFill>
      <a:ln w="38100" cap="sq">
        <a:solidFill>
          <a:srgbClr val="38B6FF"/>
        </a:solidFill>
        <a:prstDash val="solid"/>
        <a:miter/>
      </a:ln>
    </p:spPr>
  </p:sp>
  <p:sp>
    <p:nvSpPr>
      <p:cNvPr name="TextBox 4" id="4"/>
      <p:cNvSpPr txBox="true"/>
      <p:nvPr/>
    </p:nvSpPr>
    <p:spPr>
      <a:xfrm>
        <a:off x="76200" y="-9525"/>
        <a:ext cx="660400" cy="746125"/>
      </a:xfrm>
      <a:prstGeom prst="rect">
        <a:avLst/>
      </a:prstGeom>
    </p:spPr>
    <p:txBody>
      <a:bodyPr anchor="ctr" rtlCol="false" tIns="50800" lIns="50800" bIns="50800" rIns="50800">
        <a:normAutofit/>
      </a:bodyPr>
      <a:lstStyle/>
      <a:p>
        <a:pPr algn="ctr">
          <a:lnSpc><a:spcPts val="6859"/></a:lnSpc>
        </a:pPr>
      </a:p>
    </p:txBody>
  </p:sp>
</p:grpSp>"""


def create_template() -> bytes:
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)

    # Use the blank layout (index 6 in the default theme)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Remove default placeholders that come with the blank layout
    sp_tree = slide.shapes._spTree
    for ph in slide.placeholders:
        sp_tree.remove(ph._element)

    # Black background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)

    # Add the group (circle + text box)
    grp_elem = etree.fromstring(GROUP_XML)
    sp_tree.append(grp_elem)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


if __name__ == "__main__":
    output = (
        Path(__file__).parent.parent
        / "generator"
        / "templates"
        / "song_slide_template.pptx"
    )
    data = create_template()
    output.write_bytes(data)
    print(f"Template written to {output} ({len(data):,} bytes)")

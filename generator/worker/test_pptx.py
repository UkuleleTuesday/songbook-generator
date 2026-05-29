import io
from pathlib import Path

import pytest
from pptx import Presentation

from .pptx import (
    _MAX_LINES_PER_SLIDE,
    _estimate_visual_lines,
    _split_section,
    _strip_annotations,
    build_pptx,
    parse_doc_text,
)

TEMPLATE_PATH = Path(__file__).parent.parent / "templates" / "song_slide_template.pptx"

LOVE_ME_DO_TEXT = """\
Love Me Do - The Beatles

(Backing vocals)    Harmonies    148bpm    4/4    swing

[harmonica]

(G) (C) (G) (C)
(G) (C) (G) (G)

(G)Love, love me do (C)
You (G)know I love you (C)
I'll (G)always be true (C)

(G)Love, love me do (C)
You (G)know I love you (C)
"""


# --- parse_doc_text ---


def test_parse_title_extracted():
    title, _ = parse_doc_text(LOVE_ME_DO_TEXT)
    assert title == "Love Me Do"


def test_parse_metadata_skipped():
    _, sections = parse_doc_text(LOVE_ME_DO_TEXT)
    for section in sections:
        assert "bpm" not in section.lower()


def test_parse_stage_directions_included():
    _, sections = parse_doc_text(LOVE_ME_DO_TEXT)
    assert "[harmonica]" in "\n".join(sections)


def test_parse_chord_only_lines_included():
    _, sections = parse_doc_text(LOVE_ME_DO_TEXT)
    assert "(G) (C) (G) (C)" in "\n".join(sections)


def test_parse_lyric_sections_included():
    _, sections = parse_doc_text(LOVE_ME_DO_TEXT)
    assert "Love, love me do" in "\n".join(sections)


def test_parse_empty_text():
    title, sections = parse_doc_text("")
    assert title == "Untitled"
    assert sections == []


def test_parse_title_only():
    title, sections = parse_doc_text("My Song - Some Artist")
    assert title == "My Song"
    assert sections == []


def test_parse_no_artist_separator():
    title, _ = parse_doc_text("Yesterday\n\nSome lyrics here")
    assert title == "Yesterday"


def test_parse_no_annotations_strips_stage_directions():
    _, sections = parse_doc_text(LOVE_ME_DO_TEXT, include_annotations=False)
    assert "[harmonica]" not in "\n".join(sections)


def test_parse_no_annotations_keeps_chords():
    _, sections = parse_doc_text(LOVE_ME_DO_TEXT, include_annotations=False)
    assert "(G)Love, love me do" in "\n".join(sections)


def test_parse_no_annotations_drops_pure_annotation_sections():
    _, sections = parse_doc_text(
        "My Song\n\n[intro]\n\nVerse one", include_annotations=False
    )
    assert all("[intro]" not in s for s in sections)


def test_parse_no_annotations_keeps_inline_chord_lines():
    _, sections = parse_doc_text(
        "My Song\n\n[all] (G)Love me do (C)\n(G)More lyrics",
        include_annotations=False,
    )
    assert any("(G)Love me do (C)" in s for s in sections)


# --- _estimate_visual_lines ---


def test_estimate_short_section_no_wrap():
    section = "Short\nLines\nHere"
    assert _estimate_visual_lines(section) == 3


def test_estimate_wrapping_line():
    line = "x" * 60  # ceil(60/28) = 3
    assert _estimate_visual_lines(line) == 3


def test_estimate_blank_line_counts_as_one():
    assert _estimate_visual_lines("abc\n\ndef") == 3


# --- _split_section ---


def test_split_short_section_unchanged():
    section = "Line one\nLine two\nLine three"
    result = _split_section(section, _MAX_LINES_PER_SLIDE)
    assert result == [section]


def test_split_long_section_creates_multiple_chunks():
    # 5 lines each ~55 chars → ceil(55/28)=2 visual lines each → 10 total > 9
    long_line = "x" * 55
    section = "\n".join([long_line] * 5)
    result = _split_section(section, _MAX_LINES_PER_SLIDE)
    assert len(result) > 1


def test_split_slide1_budget_is_7():
    # With reserve=2, budget=7. 4 lines × ceil(55/28)=2 each = 8 > 7 → split
    long_line = "x" * 55
    section = "\n".join([long_line] * 4)
    result_full = _split_section(section, _MAX_LINES_PER_SLIDE)
    result_slide1 = _split_section(section, _MAX_LINES_PER_SLIDE - 2)
    assert len(result_slide1) >= len(result_full)


def test_split_single_very_long_line_goes_on_own_slide():
    very_long = "x" * 300  # way more than max_lines visual lines
    result = _split_section(very_long, _MAX_LINES_PER_SLIDE)
    assert len(result) == 1  # can't split a single line, keeps it as-is


def test_split_empty_section():
    assert _split_section("", _MAX_LINES_PER_SLIDE) == [""]


# --- _strip_annotations ---


def test_strip_removes_stage_direction():
    assert _strip_annotations("[harmonica]") is None


def test_strip_inline_annotation():
    assert _strip_annotations("[all] (G)Love me do (C)") == "(G)Love me do (C)"


def test_strip_mid_line_annotation():
    assert _strip_annotations("(G) (G) (G↓) [no ukes] (D7↓)") == "(G) (G) (G↓) (D7↓)"


def test_strip_preserves_chord_parens():
    assert _strip_annotations("(G)Love, love me do (C)") == "(G)Love, love me do (C)"


def test_strip_multiline_removes_empty_lines():
    assert _strip_annotations("[harmonica]\n(G)Some lyrics") == "(G)Some lyrics"


# --- build_pptx ---


def test_build_creates_valid_pptx():
    data = build_pptx("Test Song", ["Verse one\nLine two"], TEMPLATE_PATH)
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == 1


def test_build_slide_count_matches_sections():
    data = build_pptx(
        "Title", ["Section one", "Section two", "Section three"], TEMPLATE_PATH
    )
    assert len(Presentation(io.BytesIO(data)).slides) == 3


def test_build_title_appears_on_first_slide():
    data = build_pptx("My Song", ["First verse"], TEMPLATE_PATH)
    slide = Presentation(io.BytesIO(data)).slides[0]
    full_text = " ".join(
        child.text_frame.text
        for shape in slide.shapes
        if hasattr(shape, "shapes")
        for child in shape.shapes
        if child.has_text_frame
    )
    assert "My Song" in full_text


def test_build_empty_sections_produces_one_slide():
    data = build_pptx("Title", [], TEMPLATE_PATH)
    assert len(Presentation(io.BytesIO(data)).slides) == 1


def test_build_correct_slide_dimensions():
    data = build_pptx("X", ["y"], TEMPLATE_PATH)
    prs = Presentation(io.BytesIO(data))
    assert prs.slide_width.inches == pytest.approx(15.0)
    assert prs.slide_height.inches == pytest.approx(11.25)

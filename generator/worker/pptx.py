import copy
import io
import math
import re
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

from ..common.gdrive import GoogleDriveClient

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# Run properties matching the Love_Me_Do.pptx template (Open Sans, 48.99pt, white, centred)
_RUN_PROPS_NORMAL = (
    '<a:rPr xmlns:a="{ns}" lang="en-US" sz="4899">'
    '<a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>'
    '<a:latin typeface="Open Sans"/>'
    '<a:ea typeface="Open Sans"/>'
    '<a:cs typeface="Open Sans"/>'
    '<a:sym typeface="Open Sans"/>'
    "</a:rPr>"
).format(ns=_A_NS)

_RUN_PROPS_BOLD = (
    '<a:rPr xmlns:a="{ns}" lang="en-US" sz="4899" b="true">'
    '<a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>'
    '<a:latin typeface="Open Sans Bold"/>'
    '<a:ea typeface="Open Sans Bold"/>'
    '<a:cs typeface="Open Sans Bold"/>'
    '<a:sym typeface="Open Sans Bold"/>'
    "</a:rPr>"
).format(ns=_A_NS)

_PARA_PR = (
    '<a:pPr xmlns:a="{ns}" algn="ctr"><a:lnSpc><a:spcPts val="6859"/></a:lnSpc></a:pPr>'
).format(ns=_A_NS)

_METADATA_RE = re.compile(r"\bbpm\b|\b\d+/\d+\b", re.IGNORECASE)
_ANNOTATION_RE = re.compile(r"\s*\[.*?\]")

# Slide layout constants derived from template measurements:
#   text box actual width = 658pt, font = 49pt Open Sans
#   line spacing = 68.59pt exact, text box height = 743pt
_CHARS_PER_LINE = 28  # empirical chars per visual line at 49pt in 658pt box
_MAX_LINES_PER_SLIDE = 9  # floor(743/68.59)=10, minus 1 safety margin


def _strip_annotations(section: str) -> str | None:
    """Remove [stage direction] markers from a section's lines.

    Returns the cleaned section, or None if every line becomes empty.
    """
    lines = []
    for line in section.splitlines():
        cleaned = _ANNOTATION_RE.sub("", line).strip()
        if cleaned:
            lines.append(cleaned)
    return "\n".join(lines) if lines else None


def _estimate_visual_lines(section: str) -> int:
    """Estimate how many visual lines a section occupies at the template font/box size."""
    total = 0
    for line in section.splitlines():
        total += max(1, math.ceil(len(line) / _CHARS_PER_LINE)) if line.strip() else 1
    return total


def _split_section(section: str, max_lines: int) -> list[str]:
    """Greedily bin-pack lines into chunks that each fit within max_lines visual lines."""
    chunks: list[str] = []
    current: list[str] = []
    count = 0
    for line in section.splitlines():
        cost = max(1, math.ceil(len(line) / _CHARS_PER_LINE)) if line.strip() else 1
        if current and count + cost > max_lines:
            chunks.append("\n".join(current))
            current, count = [line], cost
        else:
            current.append(line)
            count += cost
    if current:
        chunks.append("\n".join(current))
    return chunks or [""]


def parse_doc_text(
    text: str, include_annotations: bool = True
) -> tuple[str, list[str]]:
    """Split a plain-text Google Doc export into (title, content_sections).

    The title is the song name (before " - Artist") from the first paragraph.
    The metadata paragraph (tempo, time sig, vocals) is always skipped.
    When include_annotations=False, [stage direction] markers are stripped and
    sections that become empty are dropped.
    """
    # Strip UTF-8 BOM and normalize Windows/old-Mac line endings to Unix
    text = text.lstrip("﻿").replace("\r\n", "\n").replace("\r", "\n")
    # Split on any sequence of blank/whitespace-only lines
    raw = [s.strip() for s in re.split(r"\n(\s*\n)+", text) if s.strip()]
    if not raw:
        return "Untitled", []

    title_line = raw[0].splitlines()[0].strip()
    title = title_line.split(" - ")[0].strip()

    content: list[str] = []
    for section in raw[1:]:
        if _METADATA_RE.search(section):
            continue
        if not include_annotations:
            section = _strip_annotations(section)
            if section is None:
                continue
        content.append(section)

    return title, content


def _make_paragraph(text: str, bold: bool = False) -> etree._Element:
    """Build an <a:p> element with a single run."""
    rpr_xml = _RUN_PROPS_BOLD if bold else _RUN_PROPS_NORMAL
    para = etree.fromstring(
        f'<a:p xmlns:a="{_A_NS}">'
        f"{_PARA_PR}"
        f"<a:r>{rpr_xml}"
        f'<a:t xml:space="preserve">{_escape(text)}</a:t>'
        f"</a:r>"
        f"</a:p>"
    )
    return para


def _make_blank_paragraph() -> etree._Element:
    return etree.fromstring(
        f'<a:p xmlns:a="{_A_NS}">'
        f"{_PARA_PR}"
        f'<a:r>{_RUN_PROPS_NORMAL}<a:t xml:space="preserve"> </a:t></a:r>'
        f"</a:p>"
    )


def _escape(text: str) -> str:
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def _find_textbox(slide) -> etree._Element | None:
    """Return the txBody of TextBox 4 inside Group 2, or None."""
    for shape in slide.shapes:
        if shape.name == "Group 2" and hasattr(shape, "shapes"):
            for child in shape.shapes:
                if child.name == "TextBox 4" and child.has_text_frame:
                    return child.text_frame._txBody
    return None


def _fill_slide(slide, title: str | None, section_text: str) -> None:
    """Replace TextBox 4 content with title (optional) + section lines."""
    txBody = _find_textbox(slide)
    if txBody is None:
        return

    # Remove existing paragraphs
    for para in txBody.findall(qn("a:p")):
        txBody.remove(para)

    if title:
        txBody.append(_make_paragraph(title, bold=True))
        txBody.append(_make_blank_paragraph())

    lines = section_text.splitlines()
    for line in lines:
        txBody.append(_make_paragraph(line))


def _copy_template_slide(prs: Presentation, template_slide) -> object:
    """Deep-copy the template slide and append it to the presentation."""
    # Copy the slide XML
    xml_copy = copy.deepcopy(template_slide._element)

    # Build a new slide part by adding a blank slide, then replace its XML
    blank_layout = template_slide.slide_layout
    new_slide = prs.slides.add_slide(blank_layout)

    # Replace the new slide's spTree and background with the template's
    new_sp_tree = new_slide.shapes._spTree
    tmpl_sp_tree = xml_copy.find(f".//{{{_P_NS}}}spTree", namespaces={"p": _P_NS})
    if tmpl_sp_tree is None:
        # Try without namespace prefix
        tmpl_sp_tree = xml_copy.find(".//{%s}spTree" % _P_NS)

    # Replace spTree contents (keep existing nvGrpSpPr and grpSpPr)
    for child in list(new_sp_tree):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag not in ("nvGrpSpPr", "grpSpPr"):
            new_sp_tree.remove(child)

    if tmpl_sp_tree is not None:
        for child in tmpl_sp_tree:
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag not in ("nvGrpSpPr", "grpSpPr"):
                new_sp_tree.append(copy.deepcopy(child))

    # Copy background
    tmpl_cSld = xml_copy.find("{%s}cSld" % _P_NS)
    if tmpl_cSld is not None:
        tmpl_bg = tmpl_cSld.find("{%s}bg" % _P_NS)
        if tmpl_bg is not None:
            new_cSld = new_slide._element.find("{%s}cSld" % _P_NS)
            if new_cSld is not None:
                existing_bg = new_cSld.find("{%s}bg" % _P_NS)
                if existing_bg is not None:
                    new_cSld.remove(existing_bg)
                new_cSld.insert(0, copy.deepcopy(tmpl_bg))

    return new_slide


def build_pptx(title: str, sections: list[str], template_path: Path) -> bytes:
    """Create a PPTX with one slide per section using the visual template."""
    prs = Presentation(str(template_path))
    template_slide = prs.slides[0]

    # Expand sections: split any that would overflow their slide
    expanded: list[str] = []
    for i, section in enumerate(sections or [""]):
        reserve = 2 if i == 0 else 0  # title + blank separator on slide 1
        expanded.extend(_split_section(section, _MAX_LINES_PER_SLIDE - reserve))
    sections = expanded

    # Populate slide 1 (already exists in the template)
    _fill_slide(template_slide, title, sections[0])

    # Add slides 2…N
    for section in sections[1:]:
        new_slide = _copy_template_slide(prs, template_slide)
        _fill_slide(new_slide, None, section)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def generate_song_pptx(
    gdrive_client: GoogleDriveClient,
    file_id: str,
    file_name: str,
    template_path: Path,
    output_path: Path,
    include_annotations: bool = True,
) -> None:
    """Download a Google Doc, parse it, and write a PPTX to output_path."""
    raw = gdrive_client.download_file(
        file_id=file_id,
        file_name=file_name,
        cache_prefix="song-pptx",
        mime_type="text/plain",
        export=True,
        use_cache=False,
    )
    text = raw.decode("utf-8")

    title, sections = parse_doc_text(text, include_annotations=include_annotations)
    data = build_pptx(title, sections, template_path)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_bytes(data)
